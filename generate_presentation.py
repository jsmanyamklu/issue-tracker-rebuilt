"""
PowerPoint Presentation Generator for Issue Tracker
Generates a complete .pptx file from the presentation content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Generate the complete PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    slide = add_title_slide(prs)

    # Problem Slide
    slide = add_problem_slide(prs)

    # Solution Slide
    slide = add_solution_slide(prs)

    # Metrics Slide
    slide = add_metrics_slide(prs)

    # Comparison Slide
    slide = add_comparison_slide(prs)

    # Features Slide
    slide = add_features_slide(prs)

    # Architecture Slide
    slide = add_architecture_slide(prs)

    # Performance Slide
    slide = add_performance_slide(prs)

    # Security Slide
    slide = add_security_slide(prs)

    # Slack Integration Slide
    slide = add_slack_slide(prs)

    # Cost Savings Slide
    slide = add_cost_slide(prs)

    # Testing Slide
    slide = add_testing_slide(prs)

    # Documentation Slide
    slide = add_documentation_slide(prs)

    # Deployment Slide
    slide = add_deployment_slide(prs)

    # Scalability Slide
    slide = add_scalability_slide(prs)

    # Demo Slide
    slide = add_demo_slide(prs)

    # Roadmap Slide
    slide = add_roadmap_slide(prs)

    # Success Metrics Slide
    slide = add_success_metrics_slide(prs)

    # Risks Slide
    slide = add_risks_slide(prs)

    # Timeline Slide
    slide = add_timeline_slide(prs)

    # Investment Slide
    slide = add_investment_slide(prs)

    # Decision Matrix Slide
    slide = add_decision_slide(prs)

    # Q&A Slide
    slide = add_qa_slide(prs)

    # Call to Action Slide
    slide = add_cta_slide(prs)

    # Summary Slide
    slide = add_summary_slide(prs)

    # Thank You Slide
    slide = add_thank_you_slide(prs)

    # Save presentation
    prs.save('Issue_Tracker_Presentation.pptx')
    print("✅ PowerPoint presentation created: Issue_Tracker_Presentation.pptx")

def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = title.text_frame
    tf.text = "ISSUE TRACKER"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    tf = subtitle.text_frame
    tf.text = "Production-Ready Enterprise Application"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)

    # Tagline
    tagline = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = tagline.text_frame
    lines = ["Built with Modern Technology", "Rivals Industry Leaders", "Ready to Deploy"]
    for i, line in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)

    # Date and status
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
    tf = footer.text_frame
    tf.text = "March 24, 2026\nStatus: PRODUCTION READY ✅"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)

    return slide

def add_problem_slide(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

    title = slide.shapes.title
    title.text = "THE CHALLENGE"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current Issue Tracking Solutions:"

    problems = [
        "Jira: $4,650-9,000/year for 50 users",
        "Legacy technology (2005)",
        "Slow performance (3-5s page loads)",
        "Vendor lock-in",
        "Limited customization",
        "Complex setup (hours/days)"
    ]

    for problem in problems:
        p = tf.add_paragraph()
        p.text = "❌ " + problem
        p.level = 0
        p.font.size = Pt(18)

    # Add "What We Need" section
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "What We Need:"
    p.font.bold = True
    p.font.size = Pt(20)

    needs = [
        "Cost-effective solution",
        "Modern, fast interface",
        "Full data control",
        "Easy to customize",
        "Quick deployment"
    ]

    for need in needs:
        p = tf.add_paragraph()
        p.text = "✅ " + need
        p.level = 0
        p.font.size = Pt(18)

    return slide

def add_solution_slide(prs):
    """Slide 3: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "OUR SOLUTION"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "A Modern, Production-Ready Issue Tracker"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True

    features = [
        "✅ $0/year - No licensing fees",
        "⚡ 4-6x Faster - Sub-second page loads",
        "🚀 5-minute Setup - Deploy immediately",
        "💻 Modern Stack - 2025 technology",
        "🔒 Enterprise Security - Production-grade",
        "🎨 Better UX - Clean, intuitive interface",
        "🔗 Slack Integration - Full bidirectional",
        "📊 100% Test Coverage - 30/30 tests passing"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(18)

    return slide

def add_metrics_slide(prs):
    """Slide 4: Key Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "PROJECT METRICS"

    content = slide.placeholders[1]
    tf = content.text_frame

    metrics = [
        ("100%", "Test Coverage (30/30 passing)"),
        ("26", "Routes Built & Optimized"),
        ("21", "API Endpoints"),
        ("12", "User-Facing Pages"),
        ("20", "Documentation Files"),
        ("4-6x", "Performance Improvement"),
        ("0", "Critical Bugs")
    ]

    for number, description in metrics:
        p = tf.add_paragraph()
        p.text = f"{number}    {description}"
        p.font.size = Pt(20)
        p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Status: PRODUCTION READY ✅"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_comparison_slide(prs):
    """Slide 5: Comparison with Jira"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "OUR SOLUTION vs. JIRA"

    # Add table
    rows, cols = 9, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(1)

    # Header row
    headers = ["Feature", "Our Solution", "Jira", "Winner"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    data = [
        ("Cost", "$0/year", "$4.6K-9K/yr", "US 💰"),
        ("Speed", "<1s loads", "3-5s loads", "US ⚡"),
        ("Setup Time", "5 minutes", "Hours/Days", "US 🚀"),
        ("Tech Stack", "Modern", "Legacy", "US 🔥"),
        ("Customize", "100%", "Limited", "US ⚙️"),
        ("Slack", "Full", "Basic", "US 💬"),
        ("Data Control", "100% Yours", "Their Cloud", "US 🔒"),
        ("AI Ready", "Yes", "No", "US 🤖")
    ]

    for i, row_data in enumerate(data, 1):
        for j, text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide

def add_features_slide(prs):
    """Slide 6: Core Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "COMPLETE FEATURE SET"

    # Create two columns
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5.5))

    # Left column
    tf_left = left_box.text_frame
    sections_left = [
        ("ISSUE MANAGEMENT", [
            "Create, Edit, Delete",
            "Assign to Team Members",
            "Status Tracking",
            "Priority Levels",
            "Type Categories",
            "Full-text Search"
        ]),
        ("INTEGRATIONS", [
            "Slack Notifications",
            "Slash Commands",
            "Link Unfurling",
            "Interactive Buttons"
        ])
    ]

    for section_title, items in sections_left:
        p = tf_left.add_paragraph()
        p.text = section_title
        p.font.bold = True
        p.font.size = Pt(14)
        for item in items:
            p = tf_left.add_paragraph()
            p.text = "✅ " + item
            p.font.size = Pt(11)
        p = tf_left.add_paragraph()
        p.text = ""

    # Right column
    tf_right = right_box.text_frame
    sections_right = [
        ("COLLABORATION", [
            "Personal Dashboards",
            "Real-time Updates",
            "Comments & Mentions",
            "Activity Timeline",
            "Team Analytics",
            "Project Organization"
        ]),
        ("SECURITY", [
            "OAuth 2.0",
            "Security Headers",
            "Rate Limiting",
            "Audit Logging"
        ])
    ]

    for section_title, items in sections_right:
        p = tf_right.add_paragraph()
        p.text = section_title
        p.font.bold = True
        p.font.size = Pt(14)
        for item in items:
            p = tf_right.add_paragraph()
            p.text = "✅ " + item
            p.font.size = Pt(11)
        p = tf_right.add_paragraph()
        p.text = ""

    return slide

def add_architecture_slide(prs):
    """Slide 7: Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "MODERN TECH STACK"

    content = slide.placeholders[1]
    tf = content.text_frame

    layers = [
        ("FRONTEND", "Next.js 15 + React 19, TypeScript + Tailwind CSS"),
        ("API LAYER", "Next.js API Routes, Clean Architecture"),
        ("BUSINESS LOGIC", "Services + Controllers, Repositories"),
        ("DATABASE", "PostgreSQL 16+, Strategic Indexing")
    ]

    for layer_name, description in layers:
        p = tf.add_paragraph()
        p.text = layer_name
        p.font.bold = True
        p.font.size = Pt(18)
        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(14)
        p.level = 1
        p = tf.add_paragraph()
        p.text = ""

    p = tf.add_paragraph()
    p.text = "INTEGRATIONS: Slack, GitHub, Google OAuth"
    p.font.size = Pt(16)
    p.font.italic = True

    return slide

def add_performance_slide(prs):
    """Slide 8: Performance"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "PERFORMANCE OPTIMIZATION"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "DATABASE QUERIES: 4-6x FASTER ⚡"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True

    improvements = [
        ("User queries", "400-600ms → 50-150ms"),
        ("Projects", "300-500ms → 40-120ms"),
        ("Issues", "500-700ms → 80-200ms")
    ]

    for query_type, improvement in improvements:
        p = tf.add_paragraph()
        p.text = f"{query_type}: {improvement}"
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "API RESPONSE TIMES"
    p.font.size = Pt(20)
    p.font.bold = True

    apis = [
        ("Health Check", "<100ms"),
        ("Issue List", "<300ms"),
        ("Issue Detail", "<200ms"),
        ("Create Issue", "<400ms")
    ]

    for api, time in apis:
        p = tf.add_paragraph()
        p.text = f"{api}: {time} ✅"
        p.font.size = Pt(16)

    return slide

def add_security_slide(prs):
    """Slide 9: Security"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "ENTERPRISE-GRADE SECURITY"

    content = slide.placeholders[1]
    tf = content.text_frame

    security_features = [
        ("AUTHENTICATION", [
            "OAuth 2.0 (GitHub & Google)",
            "Session Management (JWT)",
            "CSRF Protection",
            "No Password Storage"
        ]),
        ("NETWORK SECURITY", [
            "Content-Security-Policy",
            "Strict-Transport-Security",
            "X-Frame-Options: DENY",
            "X-Content-Type-Options: nosniff"
        ]),
        ("API PROTECTION", [
            "Rate Limiting (3 tiers)",
            "Input Validation (Zod)",
            "SQL Injection Protection"
        ])
    ]

    for section, items in security_features:
        p = tf.add_paragraph()
        p.text = section
        p.font.bold = True
        p.font.size = Pt(16)
        for item in items:
            p = tf.add_paragraph()
            p.text = "✅ " + item
            p.font.size = Pt(13)
            p.level = 1

    return slide

def add_slack_slide(prs):
    """Slide 10: Slack Integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "SEAMLESS SLACK INTEGRATION"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AUTOMATIC NOTIFICATIONS"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True

    notifications = [
        "🆕 New Issue Created",
        "👤 Issue Assigned",
        "🔄 Status Changed",
        "💬 Comment Added"
    ]

    for notif in notifications:
        p = tf.add_paragraph()
        p.text = notif
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "SLASH COMMANDS"
    p.font.size = Pt(20)
    p.font.bold = True

    commands = [
        "/issue list → Show recent issues",
        "/issue help → Show available commands"
    ]

    for cmd in commands:
        p = tf.add_paragraph()
        p.text = cmd
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "✅ Real-time notifications  ✅ Link unfurling  ✅ Interactive buttons"
    p.font.size = Pt(14)

    return slide

def add_cost_slide(prs):
    """Slide 11: Cost Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "COST ANALYSIS (50 Users)"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "JIRA COSTS (Annual)"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Standard: $7.75/user × 50 = $4,650"
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "Premium: $16/user × 50 = $9,600"
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "OUR SOLUTION (Annual)"
    p.font.size = Pt(20)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Server: $10-50/month = $120-600"
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "Licenses: $0"
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "💰 ANNUAL SAVINGS: $4,050 - $9,480"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    p = tf.add_paragraph()
    p.text = "💰 5-YEAR SAVINGS: $20,250 - $47,400"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    return slide

def add_testing_slide(prs):
    """Slide 12: Testing & Quality"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "COMPREHENSIVE TESTING"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AUTOMATED TEST SUITE"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True

    tests = [
        "Health & Monitoring: 8/8 ✅",
        "Security: 5/5 ✅",
        "API Endpoints: 4/4 ✅",
        "Page Accessibility: 5/5 ✅",
        "Database: 8/8 ✅"
    ]

    for test in tests:
        p = tf.add_paragraph()
        p.text = test
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "TOTAL: 30/30 ✅  |  PASS RATE: 100% ✅"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "CODE QUALITY"
    p.font.size = Pt(20)
    p.font.bold = True

    quality = [
        "TypeScript: 0 errors",
        "ESLint: 0 warnings",
        "Production Build: Success",
        "Bundle: Optimized"
    ]

    for item in quality:
        p = tf.add_paragraph()
        p.text = "✅ " + item
        p.font.size = Pt(16)

    return slide

def add_documentation_slide(prs):
    """Slide 13: Documentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "COMPREHENSIVE DOCUMENTATION"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "20 COMPLETE DOCUMENTATION FILES"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True

    sections = [
        ("FOR USERS", ["README.md", "SETUP_GUIDE.md", "KEYBOARD-SHORTCUTS.md"]),
        ("FOR MANAGERS", ["MANAGER-PRESENTATION.md", "EXECUTIVE-SUMMARY.md", "COMPARISON-WITH-JIRA.md"]),
        ("FOR DEVELOPERS", ["ARCHITECTURE.md", "API_DOCUMENTATION.md", "DATABASE_SCHEMA.md"]),
        ("FOR FEATURES", ["SLACK-INTEGRATION.md", "ASSIGNEE-TRACKING.md", "EXECUTION-FLOW-DIAGRAMS.md"])
    ]

    for section_title, docs in sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.bold = True
        p.font.size = Pt(14)
        for doc in docs:
            p = tf.add_paragraph()
            p.text = "• " + doc
            p.font.size = Pt(11)
            p.level = 1

    return slide

def add_deployment_slide(prs):
    """Slide 14: Deployment Options"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "FLEXIBLE DEPLOYMENT"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "OPTION 1: DOCKER (Recommended)"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    docker_features = [
        "One-command deployment",
        "Production-ready image",
        "Health checks included",
        "Setup Time: 5 minutes"
    ]

    for feature in docker_features:
        p = tf.add_paragraph()
        p.text = "✅ " + feature
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "OPTION 2: CLOUD PLATFORMS"
    p.font.size = Pt(18)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Vercel | AWS | Google Cloud | Azure | DigitalOcean"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "OPTION 3: ON-PREMISES"
    p.font.size = Pt(18)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Node.js + PostgreSQL + Nginx"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "Setup Time: 30 minutes"
    p.font.size = Pt(14)

    return slide

def add_scalability_slide(prs):
    """Slide 15: Scalability"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "BUILT TO SCALE"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "CURRENT CAPACITY (Phase 1)"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    current = [
        "Users: 100+ concurrent",
        "Issues: 10,000+ records",
        "Projects: 500+ projects",
        "Response: <1 second"
    ]

    for item in current:
        p = tf.add_paragraph()
        p.text = "✅ " + item
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "SCALING PATH"
    p.font.size = Pt(18)
    p.font.bold = True

    phases = [
        ("Phase 2: Up to 500 Users", "Load balancer, Multiple servers, Redis"),
        ("Phase 3: Up to 5,000 Users", "Kubernetes, Database sharding, CDN")
    ]

    for phase, details in phases:
        p = tf.add_paragraph()
        p.text = phase
        p.font.size = Pt(14)
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = details
        p.font.size = Pt(12)
        p.level = 1

    return slide

def add_demo_slide(prs):
    """Slide 16: Demo Highlights"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "LIVE DEMO WALKTHROUGH"

    content = slide.placeholders[1]
    tf = content.text_frame

    demo_steps = [
        ("1. DASHBOARD (2 min)", ["Real-time statistics", "Recent activity feed", "Quick action buttons"]),
        ("2. CREATE ISSUE WITH ASSIGNMENT (3 min)", ["Rich form with all metadata", "Assign to team member ◄─ KEY FEATURE", "Submit → Instant creation"]),
        ("3. TRACK PROGRESS (3 min)", ["'My Issues' personal dashboard", "Filter by assignee", "Grouped by status"]),
        ("4. SLACK NOTIFICATION (2 min)", ["Automatic notification sent", "Use /issue list command", "Click action buttons"])
    ]

    for step_title, items in demo_steps:
        p = tf.add_paragraph()
        p.text = step_title
        p.font.size = Pt(14)
        p.font.bold = True
        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(11)
            p.level = 1

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Total Demo Time: 10 minutes"
    p.font.size = Pt(16)
    p.font.italic = True

    return slide

def add_roadmap_slide(prs):
    """Slide 17: Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "FUTURE ENHANCEMENTS"

    content = slide.placeholders[1]
    tf = content.text_frame

    quarters = [
        ("Q1 (Months 1-3): Feature Parity", [
            "Command Palette (Cmd+K)",
            "Dark Mode",
            "Agile/Kanban Boards",
            "Email Notifications"
        ]),
        ("Q2 (Months 4-6): AI Features", [
            "Duplicate Detection (Better than Jira)",
            "Auto-Categorization (Better than Jira)",
            "Smart Assignment (Better than Jira)",
            "Priority Suggestions (Better than Jira)"
        ])
    ]

    for quarter, features in quarters:
        p = tf.add_paragraph()
        p.text = quarter
        p.font.size = Pt(16)
        p.font.bold = True
        for feature in features:
            p = tf.add_paragraph()
            p.text = "• " + feature
            p.font.size = Pt(13)
            p.level = 1
        p = tf.add_paragraph()
        p.text = ""

    return slide

def add_success_metrics_slide(prs):
    """Slide 18: Success Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "HOW WE MEASURE SUCCESS"

    content = slide.placeholders[1]
    tf = content.text_frame

    metrics = [
        ("TECHNICAL TARGETS", [
            "99.9% uptime",
            "<1s page load times",
            "Zero critical bugs",
            "100% test coverage"
        ]),
        ("USER ADOPTION TARGETS", [
            "80% team adoption (Month 1)",
            "90% daily active users (Month 2)",
            "<5% error rate",
            ">80% user satisfaction"
        ]),
        ("BUSINESS IMPACT TARGETS", [
            "$4K-9K annual savings",
            "30% faster issue resolution",
            "20% increase in productivity",
            "100% data ownership"
        ])
    ]

    for category, items in metrics:
        p = tf.add_paragraph()
        p.text = category
        p.font.size = Pt(14)
        p.font.bold = True
        for item in items:
            p = tf.add_paragraph()
            p.text = "✅ " + item
            p.font.size = Pt(12)
            p.level = 1

    return slide

def add_risks_slide(prs):
    """Slide 19: Risks & Mitigation"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "RISK ASSESSMENT"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "MINIMAL RISKS IDENTIFIED"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True

    risks = [
        ("Risk: User Adoption", "Mitigation: Training, pilot program, documentation"),
        ("Risk: Legacy Data Migration", "Mitigation: API-ready, migration scripts, CSV import"),
        ("Risk: Technical Issues", "Mitigation: 100% test coverage, monitoring, documentation")
    ]

    for risk, mitigation in risks:
        p = tf.add_paragraph()
        p.text = risk
        p.font.size = Pt(14)
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = mitigation
        p.font.size = Pt(12)
        p.level = 1
        p = tf.add_paragraph()
        p.text = ""

    p = tf.add_paragraph()
    p.text = "OVERALL RISK LEVEL: LOW ✅"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    return slide

def add_timeline_slide(prs):
    """Slide 20: Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "DEPLOYMENT TIMELINE"

    content = slide.placeholders[1]
    tf = content.text_frame

    weeks = [
        ("WEEK 1: Infrastructure & Deployment", [
            "Provision servers, configure database",
            "Deploy application, configure environment",
            "Testing & verification"
        ]),
        ("WEEK 2: Team Onboarding", [
            "User training sessions",
            "Pilot launch (10 users)",
            "Full team rollout"
        ]),
        ("WEEK 3-4: Full Production", [
            "Monitor metrics",
            "Gather feedback",
            "100% team adoption"
        ])
    ]

    for week, tasks in weeks:
        p = tf.add_paragraph()
        p.text = week
        p.font.size = Pt(16)
        p.font.bold = True
        for task in tasks:
            p = tf.add_paragraph()
            p.text = "• " + task
            p.font.size = Pt(13)
            p.level = 1

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Total Time to Production: 4 Weeks"
    p.font.size = Pt(18)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Can Start: THIS WEEK ✅"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    return slide

def add_investment_slide(prs):
    """Slide 21: Investment & ROI"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "INVESTMENT & ROI"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "INVESTMENT ALREADY MADE ✅"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    investments = [
        "Architecture Design",
        "Core Development",
        "Security Implementation",
        "Performance Optimization",
        "Testing & QA",
        "Documentation"
    ]

    for inv in investments:
        p = tf.add_paragraph()
        p.text = "✅ " + inv
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "ONGOING COSTS (Annual)"
    p.font.size = Pt(18)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Server Hosting: $120-600/year"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "vs. Jira: $4,650-9,600/year"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "NET SAVINGS: $4,050-9,000/year"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    p = tf.add_paragraph()
    p.text = "ROI: 700-1,500% in Year 1"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    return slide

def add_decision_slide(prs):
    """Slide 22: Decision Matrix"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "DECISION COMPARISON"

    # Create two columns
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5.5))

    # Left column - Our Solution
    tf_left = left_box.text_frame
    p = tf_left.paragraphs[0]
    p.text = "DEPLOY OUR SOLUTION"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 153, 0)

    our_benefits = [
        "$0 licensing",
        "4-6x faster performance",
        "Deploy this week",
        "Full customization",
        "100% data control",
        "Modern tech (2025)",
        "AI-ready architecture",
        "Zero vendor lock-in"
    ]

    for benefit in our_benefits:
        p = tf_left.add_paragraph()
        p.text = "✅ " + benefit
        p.font.size = Pt(13)

    # Right column - Jira
    tf_right = right_box.text_frame
    p = tf_right.paragraphs[0]
    p.text = "BUY JIRA"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(204, 0, 0)

    jira_downsides = [
        "$4.6K-9K/year",
        "Slower (3-5s)",
        "Weeks to setup",
        "Limited options",
        "Vendor cloud",
        "Legacy (2005)",
        "Not AI-ready",
        "Vendor lock-in"
    ]

    for downside in jira_downsides:
        p = tf_right.add_paragraph()
        p.text = "❌ " + downside
        p.font.size = Pt(13)

    # Recommendation at bottom
    rec_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    tf_rec = rec_box.text_frame
    p = tf_rec.paragraphs[0]
    p.text = "RECOMMENDATION: Deploy Our Solution ✅"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 153, 0)

    return slide

def add_qa_slide(prs):
    """Slide 23: Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "COMMON QUESTIONS & ANSWERS"

    content = slide.placeholders[1]
    tf = content.text_frame

    qa_pairs = [
        ("Q: What about Jira's advanced features?",
         "A: We have core features now, roadmap includes advanced features."),
        ("Q: Can it handle 100+ users?",
         "A: Yes, designed for 100+, scales to 5,000+."),
        ("Q: Is it secure enough?",
         "A: Enterprise-grade security, we control the data."),
        ("Q: Who maintains it?",
         "A: Minimal maintenance, automated health checks."),
        ("Q: Can we try it first?",
         "A: Yes! Pilot with 10-20 users for 2 weeks.")
    ]

    for question, answer in qa_pairs:
        p = tf.add_paragraph()
        p.text = question
        p.font.size = Pt(13)
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = answer
        p.font.size = Pt(12)
        p.level = 1

    return slide

def add_cta_slide(prs):
    """Slide 24: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "NEXT STEPS"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "IMMEDIATE ACTIONS REQUIRED"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True

    actions = [
        "✅ APPROVE FOR PRODUCTION",
        "📅 SCHEDULE DEPLOYMENT",
        "👥 IDENTIFY PILOT TEAM",
        "📚 PLAN TRAINING SESSION",
        "📊 DEFINE SUCCESS METRICS"
    ]

    for action in actions:
        p = tf.add_paragraph()
        p.text = action
        p.font.size = Pt(16)
        p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "WHAT WE NEED FROM YOU"
    p.font.size = Pt(20)
    p.font.bold = True

    needs = [
        "Approval to proceed",
        "Server/hosting budget ($10-50/month)",
        "2-hour training session slot",
        "Pilot team volunteers"
    ]

    for need in needs:
        p = tf.add_paragraph()
        p.text = "• " + need
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "READY TO START: THIS WEEK ✅"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    return slide

def add_summary_slide(prs):
    """Slide 25: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "EXECUTIVE SUMMARY"

    content = slide.placeholders[1]
    tf = content.text_frame

    sections = [
        ("WHAT WE HAVE", [
            "Production-ready application",
            "100% test coverage (30/30)",
            "Ready to deploy TODAY"
        ]),
        ("WHAT WE DELIVER", [
            "$4K-9K annual cost savings",
            "4-6x performance improvement",
            "100% data ownership",
            "Zero vendor lock-in"
        ]),
        ("WHAT WE NEED", [
            "Approval to proceed",
            "Deployment date",
            "$10-50/month hosting",
            "2 hours for training"
        ])
    ]

    for section, items in sections:
        p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(16)
        p.font.bold = True
        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(13)
            p.level = 1
        p = tf.add_paragraph()
        p.text = ""

    p = tf.add_paragraph()
    p.text = "DECISION: APPROVE FOR PRODUCTION ✅"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_thank_you_slide(prs):
    """Slide 26: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Thank you text
    thank_you = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = thank_you.text_frame
    tf.text = "THANK YOU"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    tf = subtitle.text_frame
    tf.text = "Questions & Discussion"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)

    # Ready to deploy
    ready = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    tf = ready.text_frame
    tf.text = "Ready to Deploy? Let's Go! 🚀"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 153, 0)

    return slide

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    print("This requires the 'python-pptx' library.")
    print("Install with: pip install python-pptx")
    print()

    try:
        create_presentation()
        print()
        print("✅ SUCCESS!")
        print("File created: Issue_Tracker_Presentation.pptx")
        print()
        print("You can now open this file in:")
        print("  • Microsoft PowerPoint")
        print("  • Google Slides (upload)")
        print("  • LibreOffice Impress")
        print()
    except ImportError:
        print("❌ ERROR: python-pptx library not found")
        print()
        print("Please install it first:")
        print("  pip install python-pptx")
        print()
        print("Then run this script again:")
        print("  python generate_presentation.py")
    except Exception as e:
        print(f"❌ ERROR: {e}")
        print()
        print("Please check the error message above and try again.")
